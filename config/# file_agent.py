# file_agent.py
import os, sys, time, hashlib, subprocess, mimetypes, pathlib, json
from datetime import datetime
from typing import Iterable, Optional, Dict, List

# Optional but recommended: pip install python-magic chardet PyPDF2 docx2txt openpyxl python-pptx
try:
    import magic  # python-magic
except Exception:
    magic = None
try:
    import chardet
except Exception:
    chardet = None
try:
    from PyPDF2 import PdfReader
except Exception:
    PdfReader = None
try:
    import docx2txt
except Exception:
    docx2txt = None
try:
    import openpyxl
except Exception:
    openpyxl = None
try:
    from pptx import Presentation
except Exception:
    Presentation = None

INDEX_DIR = "file_index"
INDEX_FILE = os.path.join(INDEX_DIR, "index.json")
CONTENT_DIR = os.path.join(INDEX_DIR, "content")

def connect():
    con = sqlite3.connect(DB)
    con.execute("PRAGMA foreign_keys=ON;")
    return con

def init_db():
    con = connect()
    for stmt in SCHEMA.strip().split(";"):
        s = stmt.strip()
        if s:
            con.execute(s)
    con.commit()
    con.close()

def sha256_file(path, buf=1024*1024):
    h = hashlib.sha256()
    with open(path, "rb") as f:
        while True:
            b = f.read(buf)
            if not b: break
            h.update(b)
    return h.hexdigest()

def detect_mime(path):
    if magic:
        try:
            return magic.from_file(path, mime=True)
        except Exception:
            pass
    mt, _ = mimetypes.guess_type(path)
    return mt or "application/octet-stream"

def read_text_best_effort(path, mime_hint):
    # cap reading to keep index light
    MAX_BYTES = 2_000_000

    # Simple handlers
    if mime_hint and mime_hint.startswith("text/"):
        with open(path, "rb") as f:
            b = f.read(MAX_BYTES)
        enc = "utf-8"
        if chardet:
            try:
                enc = chardet.detect(b)["encoding"] or "utf-8"
            except Exception:
                enc = "utf-8"
        return b.decode(enc, errors="replace")

    ext = pathlib.Path(path).suffix.lower()

    # PDF
    if ext == ".pdf" and PdfReader:
        try:
            txt = []
            reader = PdfReader(path)
            for i, page in enumerate(reader.pages[:30]):
                txt.append(page.extract_text() or "")
            return "\n".join(txt)
        except Exception:
            pass

    # DOCX
    if ext == ".docx" and docx2txt:
        try:
            return docx2txt.process(path) or ""
        except Exception:
            pass

    # XLSX
    if ext == ".xlsx" and openpyxl:
        try:
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            out = []
            for ws in wb.worksheets[:3]:
                rows = 0
                for row in ws.iter_rows(values_only=True):
                    if rows > 500: break
                    vals = [str(v) for v in row if v is not None]
                    if vals:
                        out.append(" | ".join(vals))
                        rows += 1
            return "\n".join(out)
        except Exception:
            pass

    # PPTX
    if ext == ".pptx" and Presentation:
        try:
            prs = Presentation(path)
            out = []
            for i, slide in enumerate(prs.slides[:50]):
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        out.append(shape.text)
            return "\n".join(out)
        except Exception:
            pass

    # Fallback: read first chunk as text
    try:
        with open(path, "rb") as f:
            b = f.read(MAX_BYTES)
        enc = "utf-8"
        if chardet:
            try:
                enc = chardet.detect(b)["encoding"] or "utf-8"
            except Exception:
                enc = "utf-8"
        return b.decode(enc, errors="replace")
    except Exception:
        return ""  # no preview

def upsert_file(con, path):
    try:
        st = os.stat(path)
    except FileNotFoundError:
        return

    ext = pathlib.Path(path).suffix.lower()
    mime = detect_mime(path)
    sha = sha256_file(path)
    name = os.path.basename(path)

    # Check if update needed
    cur = con.execute("SELECT size, mtime, sha256 FROM files WHERE path=?", (path,))
    row = cur.fetchone()
    if row and (row[0] == st.st_size and abs(row[1] - st.st_mtime) < 1e-6 and row[2] == sha):
        return  # unchanged

    # Extract text (best effort, capped)
    preview = read_text_best_effort(path, mime)

    con.execute("""
        INSERT INTO files(path,size,mtime,sha256,ext,mime,tags)
        VALUES(?,?,?,?,?,?,?)
        ON CONFLICT(path) DO UPDATE SET
            size=excluded.size, mtime=excluded.mtime, sha256=excluded.sha256,
            ext=excluded.ext, mime=excluded.mime
    """, (path, st.st_size, st.st_mtime, sha, ext, mime, None))

    con.execute("DELETE FROM file_fts WHERE path=?", (path,))
    con.execute("INSERT INTO file_fts(path,name,text) VALUES(?,?,?)", (path, name, preview))

def index_roots(roots: Iterable[str]):
    con = connect()
    try:
        for root in roots:
            for dirpath, _, filenames in os.walk(root):
                for fn in filenames:
                    p = os.path.join(dirpath, fn)
                    # Skip huge files if desired:
                    try:
                        if os.path.getsize(p) > 1_000_000_000:
                            continue
                    except Exception:
                        continue
                    try:
                        upsert_file(con, p)
                    except Exception as e:
                        # keep going
                        pass
        con.commit()
    finally:
        con.close()

def search(query: str, limit=20):
    con = connect()
    try:
        # Search in name and text; rank by bm25
        rows = con.execute("""
            SELECT f.path, f.ext, f.mime, f.size, f.mtime,
                   bm25(file_fts) AS rank
            FROM file_fts
            JOIN files f USING(path)
            WHERE file_fts MATCH ?
            ORDER BY rank
            LIMIT ?;
        """, (query, limit)).fetchall()
        results = []
        for path, ext, mime, size, mtime, rank in rows:
            results.append({
                "path": path,
                "ext": ext,
                "mime": mime,
                "size": size,
                "modified": datetime.fromtimestamp(mtime).isoformat(timespec="seconds"),
                "rank": rank
            })
        return results
    finally:
        con.close()

def open_file(path: str) -> Optional[str]:
    # Desktop open; return None if launched, or a message if headless
    if sys.platform.startswith("win"):
        os.startfile(path)  # type: ignore[attr-defined]
        return None
    elif sys.platform == "darwin":
        subprocess.run(["open", path], check=False)
        return None
    else:
        subprocess.run(["xdg-open", path], check=False)
        return None

if __name__ == "__main__":
    init_db()
    
    # Direct function calls without command line arguments
    # You can modify these as needed for your use case
    
    # Example: Index some directories
    # index_roots(["/path/to/directory1", "/path/to/directory2"])
    # print("Index built.")
    
    # Example: Search
    # query = "your search query"
    # for r in search(query):
    #     print(f"{r['path']}  [{r['mime']}, {r['size']}B, {r['modified']}]")
    
    # Example: Open a file
    # file_path = "/path/to/your/file"
    # msg = open_file(file_path)
    # if msg:
    #     print(msg)
